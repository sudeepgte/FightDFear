"""Generate FightDFear Product & Technical Blueprint PowerPoint (16:9)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

NAVY = RGBColor(0x12, 0x1A, 0x2F)
NAVY2 = RGBColor(0x1B, 0x2A, 0x4A)
GOLD = RGBColor(0xC9, 0xA2, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF4, 0xEE)
INK = RGBColor(0x1A, 0x1F, 0x2E)
MUTED = RGBColor(0x5C, 0x64, 0x75)
RED = RGBColor(0x8B, 0x2E, 0x2E)
AMBER = RGBColor(0x8A, 0x5A, 0x12)
GREEN = RGBColor(0x1F, 0x5C, 0x45)
ROW_ALT = RGBColor(0xEE, 0xF1, 0xF6)

W = Inches(13.333)
H = Inches(7.5)


def set_run(run, text, size=14, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def footer(slide, n, total):
    add_rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.24), Inches(8), Inches(0.24))
    p = tb.text_frame.paragraphs[0]
    set_run(p.add_run(), "FightDFear  |  Product & Technical Blueprint  |  14 Aug 2026  |  Internal", 10, False, RGBColor(0xC5, 0xCB, 0xD6))
    tb2 = slide.shapes.add_textbox(Inches(11.4), Inches(7.24), Inches(1.5), Inches(0.24))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), f"{n}  /  {total}", 10, False, GOLD)


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(0.92), NAVY)
    add_rect(slide, 0, Inches(0.92), W, Inches(0.06), GOLD)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.42))
    p = tb.text_frame.paragraphs[0]
    set_run(p.add_run(), title, 24, True, WHITE)
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.55), Inches(12.4), Inches(0.32))
        p2 = tb2.text_frame.paragraphs[0]
        set_run(p2.add_run(), subtitle, 12, False, GOLD)


def bullets(slide, items, top=1.2, left=0.5, width=12.3, size=15, color=INK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.add_run(), item, size, False, color)


def two_col(slide, left_title, left_items, right_title, right_items, top=1.15):
    add_rect(slide, Inches(0.4), Inches(top), Inches(6.05), Inches(5.85), CREAM)
    add_rect(slide, Inches(6.85), Inches(top), Inches(6.05), Inches(5.85), CREAM)
    for x, title, items in (
        (0.55, left_title, left_items),
        (7.0, right_title, right_items),
    ):
        t = slide.shapes.add_textbox(Inches(x), Inches(top + 0.12), Inches(5.75), Inches(0.4))
        set_run(t.text_frame.paragraphs[0].add_run(), title, 16, True, NAVY)
        b = slide.shapes.add_textbox(Inches(x), Inches(top + 0.55), Inches(5.75), Inches(5.1))
        tf = b.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(7)
            set_run(p.add_run(), item, 13, False, INK)


def table_slide(slide, headers, rows, left=0.35, top=1.15, width=12.65, row_h=0.36, font=11):
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(left), Inches(top), Inches(width), Inches(row_h * (len(rows) + 1)))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        set_run(p.add_run(), h, font, True, WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY2
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            col = INK
            bold = False
            text = str(val)
            if text in ("FULLY", "Ready", "P3"):
                col, bold = GREEN, True
            elif text in ("PARTIAL", "Needs work", "P1", "PLACEHOLDER"):
                col, bold = AMBER, True
            elif text in ("MISSING", "Critical", "P0", "BACKEND ONLY"):
                col, bold = RED, True
            set_run(p.add_run(), text, font, bold, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else ROW_ALT
    return tbl


def section_slide(prs, num, title, blurb):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, GOLD)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5))
    set_run(tb.text_frame.paragraphs[0].add_run(), f"SECTION  {num}", 14, True, GOLD)
    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), title, 36, True, WHITE)
    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11), Inches(1.4))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    set_run(tf3.paragraphs[0].add_run(), blurb, 16, False, RGBColor(0xC5, 0xCB, 0xD6))
    return s


def content_slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, WHITE)
    header_bar(s, title, subtitle)
    return s


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slides_meta = []  # filled as we go; footer applied after count known
    built = []

    def keep(s):
        built.append(s)
        return s

    # 1 Title
    s = keep(prs.slides.add_slide(prs.slide_layouts[6]))
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.22), H, GOLD)
    t = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.4))
    set_run(t.text_frame.paragraphs[0].add_run(), "PRODUCT  ·  TECHNICAL  ·  MARKET  BLUEPRINT", 13, True, GOLD)
    t2 = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.4))
    tf = t2.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "FightDFear", 48, True, WHITE)
    t3 = s.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.1))
    tf3 = t3.text_frame
    tf3.word_wrap = True
    set_run(tf3.paragraphs[0].add_run(), "Women Empowerment Super Platform\nComplete product blueprint, system audit & 2026 market positioning", 20, False, RGBColor(0xD8, 0xDC, 0xE6))
    t4 = s.shapes.add_textbox(Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.2))
    tf4 = t4.text_frame
    tf4.word_wrap = True
    set_run(tf4.paragraphs[0].add_run(), "Source of truth: KishorDfire repository (Spring Boot 3.4.3 · Java 17 · Flutter · JSP · MySQL 8 · Flyway V45)\nProduction: fightdfire.chethancodehub.com    ·    14 August 2026    ·    Internal — engineering, product, investors", 13, False, GOLD)

    # 2 How to read
    s = keep(content_slide(prs, "How to read this blueprint", "Youthian used only for section depth. Content is FightDFear-only."))
    table_slide(
        s,
        ["Label", "Meaning"],
        [
            ["FULLY", "Substantially implemented and usable (web and/or mobile + backend)"],
            ["PARTIAL", "Real code; important workflow, admin, payment, or client pieces missing"],
            ["BACKEND ONLY", "API/entity exists; usable Flutter (or web) path missing"],
            ["PLACEHOLDER", "Scaffold or dummy data (e.g. wallet coupons, FCM log-only)"],
            ["MISSING", "No meaningful implementation"],
            ["CURRENT / GAP / REQUIRED / FUTURE", "Never present planned work as existing"],
        ],
        row_h=0.55,
        font=13,
    )

    # 3 Exec
    s = keep(content_slide(prs, "0  ·  Executive summary", "What FightDFear is — and is not — today"))
    two_col(
        s,
        "What it is today",
        [
            "• Super platform: safety, reels, Glow, doctors, fitness, martial arts, jobs/lawyers, products, events, funding, financial literacy, creator hub",
            "• Flutter (100+ screens) + JSP web (~285 views) + REST /api/*",
            "• Real Razorpay path, JWT (mobile) + session (web), partner OTP + admin verify/reject",
            "• Admin is a first-class verification factory — not a side screen",
            "• Prod: ddl-auto=validate + Flyway (schema incidents have been real)",
        ],
        "What it is not",
        [
            "• Not ERSS-112 / police dispatch",
            "• Not a bank, UPI wallet, or SEBI-registered fund",
            "• Not an LLM copilot (no OpenAI/Gemini). Only Weka J48 safety model if ARFF exists",
            "• Not Thymeleaf (config only; templates empty)",
            "• Not a split RBAC product (no Moderator / Super Admin)",
            "• FCM tokens stored; messages not sent. SMS off by default. Wallet redeem is dummy coupons",
        ],
    )

    # 4 Snapshot
    s = keep(content_slide(prs, "2  ·  System snapshot", "From the repository, not a pitch deck"))
    table_slide(
        s,
        ["Layer", "Actual stack"],
        [
            ["Mobile", "Flutter fight_d_fear 1.0.0+1 · 102 screen files · release → production HTTPS"],
            ["Web", "Spring MVC + JSP. Thymeleaf configured, zero templates"],
            ["Admin", "AdminController /admin + extra queues (glow, martial-arts, creator, FL)"],
            ["API", "/api/auth, doctors, glow, marketplace, products, events, fitness, funding, …"],
            ["DB", "~175 entity types · MySQL 8 · Flyway V1–V45"],
            ["Auth", "BCrypt · JWT /api · session JSP · partner register-quick + email OTP"],
            ["Pay", "Razorpay create-order + verify + webhook V39 · local PAYMENT_MOCK default true"],
            ["Jobs", "ShedLock V41 · rate_limit_buckets V42"],
            ["Push", "PushNotificationService logs only — wire FCM HTTP v1"],
        ],
        row_h=0.48,
        font=12,
    )

    # 5 Roles
    s = keep(content_slide(prs, "UserType & partner types", "Do not invent Mentor / Employer / Moderator"))
    bullets(
        s,
        [
            "UserType enum: USER, ADMIN, VOLUNTEER, CENTRE, DOCTOR, STYLIST, SALON, PROVIDER, SELLER, ENTREPRENEUR, INVESTOR, FITNESS_TRAINER, EVENT_HOST",
            "Not in enum but real: FinancialEducator, DeliveryPartner, Creator (User flags), job worker (JobApplication)",
            "Lifecycle: PartnerProfileStatus (REGISTERED → … → SUSPENDED) · VerificationStatus PENDING / VERIFIED / REJECTED / CANCELLED",
            "ProviderCategory is a long catalog (lawyer, tutor, tailor, yoga, …) on ServiceProvider — not separate UserTypes",
            "ADMIN is omnipotent session admin. Spring Security does not use hasRole for UserType. CSRF is disabled. PUBLIC_URLS is large.",
        ],
        size=16,
    )

    # 6 Inventory 1
    s = keep(content_slide(prs, "3  ·  Implementation inventory (1/2)", "Code status only"))
    table_slide(
        s,
        ["Domain", "Status", "Honest note"],
        [
            ["Auth / OTP / JWT", "FULLY", "Dual session+JWT; PUBLIC_URLS sprawl"],
            ["SOS + journey + contacts (mobile)", "FULLY", "Private contacts — not 112"],
            ["Panic / recording / battery / SOS audio", "BACKEND ONLY", "JSP/API; Flutter usage not found"],
            ["Crime AI (Weka J48)", "PLACEHOLDER", "Needs crime_data.arff / safety_model.model"],
            ["Community reels (web)", "FULLY", "Moderation queue thin"],
            ["Groups / express (mobile)", "MISSING", "Web/JSP only"],
            ["Doctors / Glow / Fitness / Martial arts", "FULLY", "Core paid modules"],
            ["Jobs / lawyers / providers", "FULLY", "No Employer UserType"],
            ["Products + delivery", "PARTIAL", "Returns web-only"],
        ],
        font=12,
    )

    s = keep(content_slide(prs, "3  ·  Implementation inventory (2/2)", "Code status only"))
    table_slide(
        s,
        ["Domain", "Status", "Honest note"],
        [
            ["Events", "PARTIAL", "women_events missing on some prod DBs until V45"],
            ["Entrepreneur / investor", "PARTIAL", "Not a registered AIF — disclaimers required"],
            ["Financial literacy", "PARTIAL", "Content + booking — not regulated advice"],
            ["Creator hub", "FULLY", "Tips/cashout need admin; T&S still ops"],
            ["Wallet redeem", "PLACEHOLDER", "Dummy catalog + COUPON-timestamp"],
            ["Partner payout_balance", "PARTIAL", "Three+ money concepts"],
            ["FCM push send", "PLACEHOLDER", "Store + log only"],
            ["LLM / ChatGPT / Gemini", "MISSING", "Do not market AI copilot"],
            ["Admin partner queues", "FULLY", "Doctors, sellers, hosts, delivery, creators, educators, …"],
            ["RBAC / audit log entity", "MISSING", "P0"],
            ["iOS", "MISSING", "Android APK in repo workflow"],
        ],
        font=12,
    )

    # 7 Super platform
    s = keep(content_slide(prs, "5  ·  Super-platform map", "Document as women empowerment OS — claim only what ships"))
    table_slide(
        s,
        ["Pillar", "In code", "Do not over-claim"],
        [
            ["Safety", "SOS, journey, buddy, danger points, volunteers", "Not ERSS-112; FCM not sent; SMS off"],
            ["Community", "Reels, chat, follow, reports", "Not Instagram-scale T&S"],
            ["Jobs", "Applications, worker bookings, admin approve", "Not Naukri ATS"],
            ["Health", "Doctors, Razorpay, chat, Rx JSON", "Not hospital HIS / e-pharmacy"],
            ["Glow / fitness / skills", "Bookings, enrollment, certificates", "Not Cult.fit / university LMS"],
            ["Commerce", "Products, delivery tracking", "GST/settlement incomplete"],
            ["Funding", "Proposals, investments, meetings", "Not a SEBI fund"],
            ["AI", "Weka placeholder", "No generative AI"],
        ],
        font=12,
    )

    # 8 Current vs future
    s = keep(content_slide(prs, "CURRENT  →  GAP  →  REQUIRED  →  FUTURE", "Safety, money, admin, AI"))
    two_col(
        s,
        "P0 required (production)",
        [
            "• Flyway completeness under ddl-auto=validate (no silent Hibernate create)",
            "• SOS drill: email + SMS failover; never paywall; FCM actually send or drop the claim",
            "• Admin duty split + immutable audit log",
            "• Rotate/remove committed secrets (mail, Maps, JWT fallback, Razorpay test keys)",
            "• DPDP consent + account deletion",
            "• Legal disclaimers: health, funding, financial literacy",
            "• Replace or delete dummy wallet redeem",
        ],
        "Future (not current)",
        [
            "• ERSS-112 deep-link / dual-app guidance — not a fake dispatch network",
            "• Unified INR ledger + partner settlement reports",
            "• Moderator console + CSAM process + people",
            "• LLM safety copilot only after privacy baseline; never send SOS to third-party LLMs without consent",
            "• iOS after Android SOS/payments are proven",
            "• City densification (one metro) before national ‘super app’ marketing",
        ],
    )

    # Admin
    s = keep(content_slide(prs, "6  ·  Admin (deep)", "JSP adminDashboard — not a secondary feature"))
    table_slide(
        s,
        ["Capability", "Status", "Priority"],
        [
            ["Users approve/ban/verify/reject/delete", "FULLY", "P0"],
            ["Doctors / providers / sellers verify + request-changes", "FULLY", "P0"],
            ["Event hosts + women-events approve", "FULLY", "P0"],
            ["Salons / stylists / fitness / delivery / creators / educators", "FULLY", "P1"],
            ["Job applications approve", "FULLY", "P1"],
            ["Proposals / entrepreneurs / investors / investment release", "FULLY", "P0"],
            ["SOS list/resolve + CSV export", "FULLY", "P0"],
            ["Reported videos block + reels reward", "PARTIAL", "P0"],
            ["Broadcast + contact messages", "PARTIAL", "P1"],
            ["Product orders view / settlement UI", "PARTIAL", "P1"],
            ["Moderator vs Finance vs SuperAdmin", "MISSING", "P0"],
            ["Immutable admin audit log", "MISSING", "P0"],
        ],
        font=12,
    )

    # Mobile
    s = keep(content_slide(prs, "7  ·  Mobile blueprint", "Flutter · Navigator (no GoRouter) · 102 screen files"))
    two_col(
        s,
        "Implemented on device",
        [
            "• Auth + landing + user dashboard",
            "• SOS home, contacts, journey, buddy, danger map, reminders",
            "• Doctors (book, chat, portal, profile completion)",
            "• Glow, fitness, martial arts (user + centre + admin)",
            "• Jobs, lawyers, products, seller, delivery + live tracking",
            "• Events host, entrepreneur/investor, financial literacy, creator hub, wallet",
            "• Razorpay checkout · release host = production HTTPS",
        ],
        "Missing / weak on mobile",
        [
            "• Panic timer, recording, battery, SOS audio upload",
            "• Groups, express posts, full web chat parity",
            "• Product returns",
            "• Crash reporting / forced update / certificate pinning as product",
            "• Hindi/vernacular, iOS, proven SOS on low-network",
            "• Maps key hardcoded in maps_config.dart",
        ],
    )

    # Web
    s = keep(content_slide(prs, "8  ·  Web ecosystem", "Same MySQL · session for JSP · JWT for Flutter"))
    bullets(
        s,
        [
            "Public: /  /features  /map  /heatmap  + static beauty HTML under /index/* (marketing UI ONLY)",
            "User web: /login  /user  /qna  /chat  /video  /users/wallet  /sos  /panic  /journey  /buddy",
            "Partner JSP portals: /doctors  /salons  /stylists  /centres  /marketplace  /women-products  /entrepreneur  /investor  /creator-hub  /financial-literacy  /fitness  /women-events",
            "Admin: /admin/** plus /creator-hub/admin, /financial-literacy/admin, /qna/admin, /api/glow/admin, /api/martial-arts/admin",
            "Interaction: partners often have both a JSP dashboard and a Flutter portal against the same tables.",
        ],
        size=15,
    )

    # Safety
    s = keep(content_slide(prs, "10  ·  Safety ecosystem", "End-to-end — implemented vs planned"))
    table_slide(
        s,
        ["Flow", "Status", "Where"],
        [
            ["SOS trigger / status / cancel / history", "FULLY", "API + Flutter + JSP"],
            ["Notify trusted contacts (email)", "FULLY", "SosAsyncNotificationService"],
            ["SMS (Twilio/Msg91)", "PARTIAL", "sms.enabled=false default"],
            ["FCM to device", "PLACEHOLDER", "Log only"],
            ["Admin SOS resolve", "FULLY", "/admin/sos"],
            ["Journey + ShedLock escalation", "FULLY", "/api/journey"],
            ["Buddy mode", "PARTIAL", "API + Flutter + JSP"],
            ["Panic timer / live location page / recording", "BACKEND ONLY", "JSP"],
            ["112 / police dispatch", "MISSING", "Must deep-link, not fake"],
        ],
        font=12,
    )

    # Community jobs health etc condensed
    s = keep(content_slide(prs, "11–18  ·  Vertical modules", "Status from controllers + Flutter screens"))
    table_slide(
        s,
        ["Module", "Users", "Status", "Production caveat"],
        [
            ["Community / reels", "USER", "PARTIAL", "Mobile groups missing; need moderators"],
            ["Jobs", "Worker + admin", "FULLY", "No Employer role / ATS"],
            ["Martial arts + FL", "Centre / educator", "FULLY", "FL is education, not advice"],
            ["Doctors", "USER + DOCTOR", "FULLY", "Clinical/legal SOP"],
            ["Glow", "USER + SALON/STYLIST", "FULLY", "Settlement"],
            ["Products", "USER + SELLER + delivery", "PARTIAL", "Returns web-only"],
            ["Lawyer / fitness / events", "Providers", "FULLY", "Events schema ops (V45)"],
            ["Funding", "Entrepreneur / investor", "PARTIAL", "Securities disclaimer"],
            ["Creator hub", "Creator", "FULLY", "Cashout admin + T&S"],
        ],
        font=12,
    )

    # Wallet
    s = keep(content_slide(prs, "19  ·  Wallet, coins, payouts", "Three overlapping money systems"))
    table_slide(
        s,
        ["Mechanism", "Where", "Status"],
        [
            ["rewardPoints + WalletTransaction", "User + wallet APIs", "Real points; dummy rewards list"],
            ["Redeem catalog", "WalletController + MobileWalletController", "PLACEHOLDER coupons"],
            ["walletBalance (Double)", "User", "Incomplete INR wallet product"],
            ["Razorpay orders", "Doctors, glow, enrollment, products", "Real path; mock flag"],
            ["payout_balance on partners", "Doctor, salon, seller, …", "Ops incomplete"],
            ["Creator tips / cashout", "creator_* tables", "Admin approve"],
        ],
        font=13,
    )
    note = s.shapes.add_textbox(Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.5))
    set_run(note.text_frame.paragraphs[0].add_run(), "Rule: SOS / panic must never be paywalled.", 14, True, RED)

    # AI
    s = keep(content_slide(prs, "20  ·  AI", "Do not advertise what is not wired"))
    bullets(
        s,
        [
            "LLM (OpenAI / Gemini / ChatGPT): MISSING in Java and Dart.",
            "AISafetyService: Weka J48. Loads safety_model.model or trains from crime_data.arff in project root. If files absent → predictions unavailable. Flutter crime-check not found.",
            "SafetyCheckController /safety/check exists on web — treat as experimental.",
            "Future: moderation assist and route risk only after DPDP, with human-in-the-loop. Never send SOS/location to third-party LLMs without explicit consent.",
        ],
        size=16,
    )

    # RBAC
    s = keep(content_slide(prs, "21  ·  Permissions (actual roles)", "ADMIN sees everything; no moderator split"))
    table_slide(
        s,
        ["Module", "USER", "Partner*", "ADMIN"],
        [
            ["Own profile", "VCU", "VCU", "Approve / delete"],
            ["SOS", "Create", "Volunteer respond", "Resolve"],
            ["Reels", "VCUM", "Creator extras", "Moderate / block"],
            ["Doctor / Glow / Fitness book", "Create", "Manage supply", "Verify partner"],
            ["Jobs / products", "Apply / order", "Worker / seller", "Approve"],
            ["Events / funding", "Register / view", "Host / entrepreneur / investor", "Approve + release $"],
        ],
        font=13,
    )
    n = s.shapes.add_textbox(Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.5))
    set_run(n.text_frame.paragraphs[0].add_run(), "*Partner = CENTRE, DOCTOR, SALON/STYLIST, PROVIDER, SELLER, EVENT_HOST, FITNESS_TRAINER, ENTREPRENEUR, INVESTOR, plus educator/delivery/creator tables.", 11, False, MUTED)

    # Journeys
    s = keep(content_slide(prs, "22  ·  End-to-end journeys (as coded)", "Not aspirational hiring funnels"))
    bullets(
        s,
        [
            "New user: Flutter /api/auth → /api/landing → dashboard. Web: /login session.",
            "Emergency: SOS APIs → email (SMS if enabled) → admin /sos → resolve. No automatic police.",
            "Job seeker: jobs register-quick + OTP → apply → admin may approve → worker bookings.",
            "Provider (doctor/salon/lawyer/trainer/centre): OTP → profile completion → partner_profile_status → admin verify/reject/request-changes → bookings.",
            "Seller: seller auth → catalog → orders → delivery partner → tracking. Returns on web.",
            "Entrepreneur: register → proposal → admin approve → investor interest → chat/meetings.",
        ],
        size=16,
    )

    # Architecture
    s = keep(content_slide(prs, "23–25  ·  Architecture, data, APIs", "Monolith WAR · mixed auth · DB queries not Elasticsearch"))
    two_col(
        s,
        "As built",
        [
            "• Flutter (JWT) + JSP (session) + Admin session → Spring Boot 3.4 :8084",
            "• MySQL 8 · Flyway V1–V45 · local ddl-auto=update · prod validate",
            "• FlywayRepairConfig: repair() then migrate() on startup (can hide checksum issues)",
            "• Razorpay · SMTP (Gmail in properties) · Google Maps · WebSocket/STOMP",
            "• Entity clusters: identity, safety, community, glow, health, fitness, martial, FL, marketplace, products, events, funding, creator, payments",
        ],
        "Proposed (not existing)",
        [
            "• admin_audit_log · moderation_case · unified ledger_entry · consent_record",
            "• /api/v1 versioning + OpenAPI",
            "• Idempotent SOS + 112 webhook",
            "• Role matchers instead of PUBLIC_URLS sprawl",
            "• Object storage ACL for ID docs (see STORAGE_MIGRATION.md)",
        ],
    )

    # Security
    s = keep(content_slide(prs, "26  ·  Security & privacy", "Visible from code — treat as P0"))
    table_slide(
        s,
        ["Finding", "Evidence", "Priority"],
        [
            ["Committed secrets", "application.properties: mail, Razorpay test, Maps, JWT default", "P0"],
            ["Hardcoded JWT fallback", "JwtUtil if env blank", "P0"],
            ["Maps key in Flutter", "maps_config.dart + GoogleMapsService.FALLBACK_KEY", "P0"],
            ["CSRF disabled", "SecurityConfig", "P0"],
            ["No hasRole matchers", "authenticated() after huge permitAll list", "P0"],
            ["FCM not sent / SMS off", "PushNotificationService · sms.enabled=false", "P0"],
            ["Local ddl-auto=update", "Drift vs Flyway", "P0"],
            ["DPDP erasure / consent UX", "Admin delete user only", "P0"],
            ["ID docs / PHI", "Paths + MedicalDetails, Rx JSON — extra sensitivity", "P1"],
        ],
        font=12,
    )

    # Market
    s = keep(content_slide(prs, "27  ·  2026 market (sourced)", "Registered crime ≠ prevalence. Vendor TAM ≠ official."))
    bullets(
        s,
        [
            "NCRB Crime in India 2024: 4,41,534 registered crimes against women (vs 4,48,211 in 2023, −1.4%). Rate 64.6 per lakh women (2024) vs 66.2 (2023). Sources: NCRB via The Print / CEDA Ashoka (2026). Under-reporting remains.",
            "PIB Mar 2025: Nirbhaya Fund allocation ₹7,712.85 crore through FY 2024–25; utilisation ~₹5,846 crore (~76%). ERSS-112 in all 36 States/UTs. WHL-181: 1.95 crore+ calls, 81.64 lakh+ women assisted (PIB figures).",
            "Table stakes: 112 India app, Himmat Plus, Safetipin. FightDFear must deep-link 112 — not replace it.",
            "Livelihood: Sitha (company-stated 1M women entrepreneurs by 2027; SheJobs ~1 lakh DB — not independently audited). SHE-MART (site: 60,000+ products). Kaabil (Mahindra) / HerJobs for careers.",
            "Safety-app TAM reports (Intel Market Research et al.) are vendor estimates — do not cite as audited.",
        ],
        size=14,
    )

    # Competitors
    s = keep(content_slide(prs, "28  ·  Competitor landscape", "Closest overlaps — not a copied Youthian list"))
    table_slide(
        s,
        ["Category", "Names", "Lesson"],
        [
            ["Official safety", "112 India, Himmat Plus, WHL-181", "Trust + dispatch; free; do not compete"],
            ["Safety apps", "Safetipin, bSafe, Life360, Noonlight, Raksha", "SOS polish vs India livelihood"],
            ["Women gigs/commerce", "Sitha, SHE-MART, Womaniya on GeM", "Closest livelihood overlap"],
            ["Jobs", "Kaabil, HerJobs, Naukri (indirect)", "Distribution threat"],
            ["Health / beauty / social", "Practo, Instagram (indirect)", "Do not fight on their home turf"],
        ],
        row_h=0.55,
        font=13,
    )

    s = keep(content_slide(prs, "29  ·  Feature comparison", "FightDFear column = code status"))
    table_slide(
        s,
        ["Feature", "FightDFear", "112", "Life360", "Sitha", "Kaabil"],
        [
            ["SOS / dispatch", "Private contacts", "Official", "Circle", "—", "—"],
            ["Journey / location", "PARTIAL", "Emergency", "Strong", "Job track", "—"],
            ["Jobs / gigs", "FULLY core", "—", "—", "Strong", "Strong"],
            ["Doctors / beauty", "FULLY core", "—", "—", "Wellness gigs", "—"],
            ["Products", "PARTIAL", "—", "—", "Handmade", "—"],
            ["Community video", "FULLY web", "—", "—", "—", "—"],
            ["AI / LLM", "MISSING", "—", "Some", "Matching claims", "Matching"],
            ["Verification ops", "Admin queues", "Govt ID", "Family", "KYC claimed", "Employer"],
        ],
        font=12,
    )

    # Pricing
    s = keep(content_slide(prs, "30  ·  Competitor pricing", "Do not invent undisclosed commissions"))
    bullets(
        s,
        [
            "112 / Himmat Plus / 181: free (public funds / Nirbhaya).",
            "Life360 official US pages (2026): Silver $7.99 · Gold $14.99 · Platinum $24.99 / month. Some regional pages show $9.99 Silver — verify geo. Not India list price.",
            "Noonlight: core panic free (US); Premium ~$9.99/mo reported (ImAlive Jul 2026).",
            "Hollie Guard Extra: £7.99/mo (ImAlive Jul 2026, verified in that comparison).",
            "Sitha: bookings + booster plans (site). Commission not publicly disclosed.",
            "SHE-MART: DPIIT sellers 0% commission 6 months (site claim).",
            "FightDFear code: commission_percent / platform_fee fields exist. Consumer subscription SKU is not a complete product.",
        ],
        size=15,
    )

    # Gaps + UVP
    s = keep(content_slide(prs, "31–32  ·  Market gaps & differentiation", "Why one app vs ten — if trust holds"))
    two_col(
        s,
        "Genuine gaps",
        [
            "• 112 is trusted but not a livelihood OS",
            "• Sitha / SHE-MART / Kaabil do not own the emergency graph",
            "• Safety-only apps churn after install — they don’t pay the bills",
            "• Verification theatre: queues without ops people fail",
            "• Do not pretend to replace 112 or Instagram or UPI",
        ],
        "Honest UVP",
        [
            "• Verified women-centric services PLUS personal emergency tools",
            "• Shared identity: reputation can travel from doctor booking to worker hire",
            "• Network effects only city-by-city (supply density)",
            "• Not: ‘AI police’. Not: a bank. Not: a national super-app on day one",
        ],
    )

    # Monetization
    s = keep(content_slide(prs, "33  ·  Monetization (SOS-safe)", "Take-rate on services; never on panic"))
    table_slide(
        s,
        ["Stream", "Fit", "Caution"],
        [
            ["Glow / products / jobs commission", "High", "Transparent fees"],
            ["Doctor / fitness take-rate", "High", "Already in schema"],
            ["Event tickets", "Medium", "After events table stable"],
            ["Partner boosters / subscriptions", "Medium", "Sitha-like"],
            ["CSR / NGO / Safe City", "Strategic", "Long procurement"],
            ["Premium AI", "Later", "After trust + actual LLM"],
            ["SOS / panic", "Never paywall", "Ethics + incidents"],
        ],
        font=13,
    )

    # Roadmap
    s = keep(content_slide(prs, "34  ·  Product roadmap", "Prioritize safety, schema, trust, then growth"))
    table_slide(
        s,
        ["Horizon", "Work"],
        [
            ["NOW (0–90 days)", "Validate-green Flyway · SOS drill · secrets · audit log · dummy wallet · DPDP · health/funding disclaimers"],
            ["NEXT (3–6 months)", "Unified ledger · moderator console · 112 deep-link · vernacular · FCM real send · load tests"],
            ["LATER (6–12 months)", "One-metro density · iOS decision · careful moderation-assist AI"],
            ["FUTURE (12–24 months)", "CSR / government · Safetipin-class data partnership · still not a bank or AIF"],
        ],
        row_h=0.7,
        font=14,
    )

    # Production
    s = keep(content_slide(prs, "35  ·  Production readiness", "Green / amber / red from repo evidence"))
    table_slide(
        s,
        ["Area", "Status", "Note"],
        [
            ["Functionality breadth", "Needs work", "Wide, uneven"],
            ["Database / Flyway", "Critical", "validate vs Hibernate-era tables"],
            ["Security / DPDP / secrets", "Critical", "CSRF off, PUBLIC_URLS, committed keys"],
            ["Testing", "Critical", "Few tests vs ~108 controllers"],
            ["Payments", "Needs work", "Real path; mock must stay false in prod"],
            ["Push / SMS", "Critical", "FCM not sent; SMS off"],
            ["Admin / moderation", "Needs work", "Queues exist; roles and people do not"],
            ["Monitoring", "Needs work", "Actuator/prometheus; APM thin"],
        ],
        font=13,
    )

    # Master gap
    s = keep(content_slide(prs, "36  ·  Master gap table", "P0 = do not scale marketing until moved"))
    table_slide(
        s,
        ["Module", "Status", "Priority"],
        [
            ["Schema / Flyway", "PARTIAL", "P0"],
            ["SOS reliability + 112 guidance", "PARTIAL", "P0"],
            ["Secrets / CSRF / RBAC / audit", "MISSING", "P0"],
            ["Community T&S staffing", "PARTIAL", "P0"],
            ["Funding / health legal perimeter", "PARTIAL", "P0"],
            ["Wallet honesty", "PLACEHOLDER", "P1"],
            ["FCM send", "PLACEHOLDER", "P1"],
            ["Events related tables", "PARTIAL", "P1"],
            ["LLM AI", "MISSING", "P3"],
            ["iOS", "MISSING", "P2"],
        ],
        font=13,
    )

    # Close
    s = keep(content_slide(prs, "37  ·  Strategic close", "For a new team in 90 / 6 / 24 months"))
    two_col(
        s,
        "Today vs should become",
        [
            "TODAY: Broad women-centric Super App MVP in production clothing — real APIs and admin queues, fighting schema/ops and Trust & Safety maturity.",
            "SHOULD BECOME: Trusted OS for a woman’s safety graph and verified local economy in selected Indian cities. SOS free forever; services fund the company.",
            "WEAKNESSES: secrets, CSRF, Flyway, money sprawl, god-admin, no LLM (don’t advertise it), no police integration, test debt.",
        ],
        "Recommended clock",
        [
            "90 DAYS: schema green, SOS drill, audit log, wallet honesty, DPDP, one-city supply quality, rotate secrets.",
            "6 MONTHS: ledger, moderator, 112 deep-link, vernacular, real push, iOS decision.",
            "12–24 MONTHS: partnerships, careful AI, still not a bank.",
            "BUSINESS RISKS: winning every vertical vs specialists; funding-module perception; a real SOS miss; Play permission policy.",
        ],
    )

    # Appendix
    s = keep(content_slide(prs, "Appendix  ·  How to use with engineering", "Companion documents in the repo"))
    bullets(
        s,
        [
            "Full narrative: docs/FightDFear_Product_Technical_Blueprint.md",
            "Ops: docs/DEPLOYMENT.md · PRODUCTION_READINESS_REPORT.md · STORAGE_MIGRATION.md · REALTIME_MULTI_INSTANCE.md",
            "QA: MASTER_QA_TEST_PLAN.md and QA_*_MODULE.md per vertical",
            "This deck is implementation-true as of 14 Aug 2026. If code and slides disagree later, trust the repository.",
            "Youthian_blueprint.pptx was a 38-slide image deck used only as a structural depth reference — no Youthian product content was copied.",
        ],
        size=16,
    )

    total = len(built)
    for i, sl in enumerate(built, 1):
        # skip full-bleed title and section? still footer on title might clash
        if i == 1:
            tb = sl.shapes.add_textbox(Inches(0.9), Inches(6.9), Inches(11), Inches(0.3))
            set_run(tb.text_frame.paragraphs[0].add_run(), f"1  /  {total}", 11, False, GOLD)
            continue
        footer(sl, i, total)

    out = Path(__file__).resolve().parent / "FightDFear_Product_Technical_Blueprint.pptx"
    prs.save(str(out))
    print(f"Wrote {out} ({total} slides)")
    return out


if __name__ == "__main__":
    build()
