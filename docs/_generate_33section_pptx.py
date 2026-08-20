"""FightDFear 33-section blueprint PPT — brand from Flutter/JSP: rose #F43F5E + navy #1E1B4B."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Brand (user_dashboard_screen.dart / JSP CSS)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
NAVY = RGBColor(0x1E, 0x1B, 0x4B)
NAVY2 = RGBColor(0x31, 0x2E, 0x81)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1E, 0x1B, 0x4B)
MUTED = RGBColor(0x64, 0x73, 0x8B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
SOS = RGBColor(0xFF, 0xF1, 0xF2)
OK = RGBColor(0x16, 0xA3, 0x4A)
WARN = RGBColor(0xC2, 0x41, 0x0C)
TEAL = RGBColor(0x20, 0xC9, 0x97)
ALT = RGBColor(0xFF, 0xFC, 0xFD)
ROSE_DK = RGBColor(0xBE, 0x12, 0x3C)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Calibri"


def run(p, text, size=14, bold=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def rect(slide, l, t, w, h, fill, line=None, radius=False):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    if radius:
        try:
            sh.adjustments[0] = 0.08
        except Exception:
            pass
    return sh


def tb(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run(p, text, size, bold, color)
    return box


def footer(slide, n, total):
    rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY)
    rect(slide, 0, Inches(7.22), Inches(0.12), Inches(0.28), ROSE)
    tb(slide, Inches(0.35), Inches(7.235), Inches(10), Inches(0.24),
       "FightDFear  ·  Product · Technical · Market · Business Blueprint  ·  14 Aug 2026  ·  Internal",
       10, False, RGBColor(0xCB, 0xD5, 0xE1))
    t = slide.shapes.add_textbox(Inches(11.55), Inches(7.235), Inches(1.5), Inches(0.24))
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run(t.text_frame.paragraphs[0], f"{n} / {total}", 10, True, ROSE)


def header(slide, kicker, title, subtitle=None):
    rect(slide, 0, 0, W, H, BG)
    rect(slide, 0, 0, W, Inches(0.98), NAVY)
    rect(slide, 0, Inches(0.98), W, Inches(0.07), ROSE)
    tb(slide, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.22), kicker, 10, True, ROSE)
    tb(slide, Inches(0.4), Inches(0.30), Inches(12.5), Inches(0.38), title, 22, True, WHITE)
    if subtitle:
        tb(slide, Inches(0.4), Inches(0.66), Inches(12.5), Inches(0.28), subtitle, 11, False, RGBColor(0xCB, 0xD5, 0xE1))


def content(prs, kicker, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, kicker, title, subtitle)
    return s


def divider(prs, num, title, blurb):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, ROSE)
    rect(s, 0, Inches(7.22), W, Inches(0.28), ROSE)
    tb(s, Inches(0.7), Inches(2.35), Inches(11.5), Inches(0.35), f"SECTION  {num:02d}  OF  33", 13, True, ROSE)
    tb(s, Inches(0.7), Inches(2.8), Inches(11.8), Inches(1.3), title, 32, True, WHITE)
    box = s.shapes.add_textbox(Inches(0.7), Inches(4.25), Inches(11.2), Inches(1.6))
    box.text_frame.word_wrap = True
    run(box.text_frame.paragraphs[0], blurb, 16, False, RGBColor(0xCB, 0xD5, 0xE1))
    return s


def bullets(slide, items, top=1.22, left=0.45, width=12.4, size=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run(p, item, size, False, INK)


def card(slide, l, t, w, h, title, body, accent=ROSE):
    sh = rect(slide, Inches(l), Inches(t), Inches(w), Inches(h), WHITE, BORDER, True)
    rect(slide, Inches(l), Inches(t), Inches(0.08), Inches(h), accent)
    tb(slide, Inches(l + 0.22), Inches(t + 0.12), Inches(w - 0.35), Inches(0.32), title, 13, True, NAVY)
    box = slide.shapes.add_textbox(Inches(l + 0.22), Inches(t + 0.46), Inches(w - 0.35), Inches(h - 0.58))
    box.text_frame.word_wrap = True
    run(box.text_frame.paragraphs[0], body, 12, False, MUTED)
    return sh


def table(slide, headers, rows, left=0.35, top=1.18, width=12.65, rh=0.38, fs=11):
    shp = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(rh * (len(rows) + 1)))
    tbl = shp.table
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = ""
        p = c.text_frame.paragraphs[0]
        run(p, h, fs, True, WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    status = {
        "FULLY": OK, "PARTIAL": WARN, "PLACEHOLDER": WARN, "MISSING": ROSE_DK,
        "BACKEND ONLY": ROSE_DK, "UI ONLY": WARN, "P0": ROSE_DK, "P1": WARN,
        "P2": TEAL, "P3": MUTED, "FACT": NAVY, "RECOMMENDATION": ROSE,
        "FUTURE": MUTED, "CURRENT": OK, "Verified": OK, "Estimated": WARN,
        "Company-stated": WARN, "Third-party": MUTED,
    }
    for r, row in enumerate(rows, 1):
        for i, val in enumerate(row):
            c = tbl.cell(r, i)
            c.text = ""
            p = c.text_frame.paragraphs[0]
            text = str(val)
            col = status.get(text, INK)
            run(p, text, fs, text in status, col)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xFF, 0xF1, 0xF2)
    return tbl


def flow_row(slide, labels, top=2.4):
    n = len(labels)
    gap = 0.12
    w = (12.4 - gap * (n - 1)) / n
    x = 0.45
    for i, lab in enumerate(labels):
        rect(slide, Inches(x), Inches(top), Inches(w), Inches(0.72), NAVY if i == 0 or i == n - 1 else WHITE, ROSE if i not in (0, n - 1) else None, True)
        tb(slide, Inches(x), Inches(top + 0.18), Inches(w), Inches(0.42), lab, 11, True,
           WHITE if i == 0 or i == n - 1 else NAVY, PP_ALIGN.CENTER)
        if i < n - 1:
            tb(slide, Inches(x + w - 0.05), Inches(top + 0.18), Inches(0.22), Inches(0.4), ">", 16, True, ROSE, PP_ALIGN.CENTER)
        x += w + gap


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    built = []

    def keep(s):
        built.append(s)
        return s

    # COVER
    s = keep(prs.slides.add_slide(prs.slide_layouts[6]))
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.18), H, ROSE)
    rect(s, 0, Inches(6.55), W, Inches(0.95), ROSE)
    tb(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.3), "PRODUCT  ·  TECHNICAL  ·  MARKET  ·  BUSINESS  BLUEPRINT", 12, True, ROSE)
    tb(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.1), "FightDFear", 54, True, WHITE)
    tb(s, Inches(0.7), Inches(3.15), Inches(11.5), Inches(1.1),
       "Women Empowerment Super Platform\n33-section enterprise blueprint  ·  Codebase as source of truth  ·  2026 market research",
       18, False, RGBColor(0xE2, 0xE8, 0xF0))
    tb(s, Inches(0.7), Inches(4.55), Inches(11.5), Inches(1.4),
       "Stack: Spring Boot 3.4.3 · Java 17 · Flutter · JSP · MySQL 8 · Flyway V45 · Razorpay\n"
       "Production: fightdfire.chethancodehub.com    ·    14 August 2026    ·    Internal\n"
       "Youthian PPT used only as structural template. FightDFear content, brand (#F43F5E / #1E1B4B) and facts are original.",
       13, False, RGBColor(0xCB, 0xD5, 0xE1))
    tb(s, Inches(0.7), Inches(6.72), Inches(12), Inches(0.5),
       "Brand palette from mobile/lib/screens/user/user_dashboard_screen.dart and JSP CSS  ·  Not a safety-only app",
       12, True, WHITE)

    # How to read
    s = keep(content(prs, "READING GUIDE", "Fact vs recommendation", "Do not mix categories. Youthian = structure only."))
    table(s, ["Label", "Meaning"], [
        ["FACT", "Verified from FightDFear code or a cited external source"],
        ["CURRENT IMPLEMENTATION", "Exists in the KishorDfire repository today"],
        ["MARKET RESEARCH", "2026 external source with date — not invented"],
        ["INFERENCE", "Reasonable conclusion from evidence"],
        ["RECOMMENDATION", "What FightDFear should do"],
        ["FUTURE", "Not implemented — do not market as live"],
        ["FULLY / PARTIAL / BACKEND ONLY / PLACEHOLDER / MISSING", "Implementation status from audit"],
    ], rh=0.58, fs=13)

    # 01
    keep(divider(prs, 1, "Introduction", "What FightDFear is, why it exists, who it serves."))
    s = keep(content(prs, "01  INTRODUCTION", "Women Empowerment Super Platform", "FACT: category claim is product strategy. CURRENT: safety is one of many modules."))
    card(s, 0.4, 1.22, 6.1, 1.7, "What it is",
         "A women-centric operating system combining personal emergency tools with verified livelihood, health, skills, commerce, community and funding — one identity, one admin trust layer.", NAVY)
    card(s, 6.7, 1.22, 6.2, 1.7, "Why it exists",
         "Women in India still juggle 112 + job apps + clinics + Instagram + marketplaces. Fragmentation is the product problem. Safety without opportunity does not retain. Opportunity without safety does not trust.", ROSE)
    card(s, 0.4, 3.1, 4.05, 2.0, "Who it serves", "Women users; doctors, salons, workers, sellers, trainers, centres, hosts, educators, creators, entrepreneurs, investors, delivery partners; operators (ADMIN).", TEAL)
    card(s, 4.6, 3.1, 4.05, 2.0, "Vision", "A trusted city-level OS where a woman can call for help and also earn, learn, heal and sell — without five logins and unverified strangers.", NAVY2)
    card(s, 8.8, 3.1, 4.1, 2.0, "Mission", "Ship reliable SOS (never paywalled), verified partners, and honest status. Deep-link 112. Earn on services, not on panic.", ROSE)
    tb(s, Inches(0.45), Inches(5.3), Inches(12.4), Inches(1.6),
       "FACT (NCRB 2024, reported 2026): 4,41,534 registered crimes against women; rate 64.6 per lakh. Registered crime ≠ true prevalence.\n"
       "FACT (MoSPI PLFS): female LFPR 40.0% usual status CY 2025; 32.7% CWS June 2026. Urban female CWS LFPR 24.8% (June 2026 bulletin).\n"
       "INFERENCE: the economic + safety gap is simultaneous. A safety-only app leaves urban women under-employed; a jobs-only app leaves the emergency graph empty.",
       13, False, MUTED)

    # 02
    keep(divider(prs, 2, "Project overview", "Ecosystem map — only CURRENT modules are marked live."))
    s = keep(content(prs, "02  PROJECT OVERVIEW", "Ecosystem — live vs planned", "CURRENT IMPLEMENTATION from controllers + Flutter screens"))
    table(s, ["Pillar", "Status", "In the repo", "Do not claim"], [
        ["Safety / SOS / journey / buddy", "PARTIAL", "SOS+journey+contacts mobile FULLY; panic/audio/recording web", "Not ERSS-112"],
        ["Community / video / chat", "PARTIAL", "JSP reels FULLY; mobile groups MISSING", "Not Instagram"],
        ["Jobs / lawyers / marketplace services", "FULLY", "Worker bookings, ProviderCategory incl. WOMEN_LAWYER", "Not Naukri ATS"],
        ["Glow / doctors / fitness / martial arts", "FULLY", "Razorpay + portals + admin verify", "Not Practo/Cult.fit scale"],
        ["Products + delivery", "PARTIAL", "Returns web-only", "Not national logistics"],
        ["Events / funding / FL / creator", "PARTIAL", "V45 women_events; FL = education not banking", "Not a SEBI fund / bank"],
        ["Wallet / coins", "PLACEHOLDER", "Dummy redeem catalog both web + /api/wallet", "Not a PPI license"],
        ["AI / LLM", "MISSING", "Only Weka J48 if ARFF present", "No ChatGPT copilot"],
        ["Admin verification factory", "FULLY", "Queues for doctors, sellers, hosts, delivery, creators, educators…", "No Moderator role"],
    ], fs=11)

    s = keep(content(prs, "02  PROJECT OVERVIEW", "Surfaces that actually exist", "FACT: Thymeleaf is configured; templates folder is empty. UI is JSP + Flutter."))
    card(s, 0.4, 1.22, 4.05, 2.35, "Mobile  ·  Flutter", "102 screen files. JWT. Release host = production HTTPS. SOS, doctors, Glow, jobs, products, events, funding, FL, creator, wallet, partner profile-completion.", ROSE)
    card(s, 4.6, 1.22, 4.05, 2.35, "Web  ·  JSP (~285 views)", "Public /features /map /heatmap. User SOS/video/chat/qna. Partner dashboards: doctors, salons, centres, marketplace, products, entrepreneur, investor, creator, FL.", NAVY)
    card(s, 8.8, 1.22, 4.1, 2.35, "Admin  ·  /admin + extras", "God-admin session. Verify/reject/request-changes. SOS resolve. Broadcast. Reports CSV. Glow/martial-arts/creator/FL extra admin UIs.", TEAL)
    card(s, 0.4, 3.75, 6.1, 2.2, "API", "/api/auth, landing, sos, journey, buddy, videos, wallet, doctors, glow, marketplace, jobs, products, delivery, martial-arts, fitness, women-events, entrepreneur, investor, financial-literacy, creator-hub, chat.", NAVY2)
    card(s, 6.7, 3.75, 6.2, 2.2, "Data / pay / jobs", "MySQL 8 · ~175 entity types · Flyway V1–V45. Razorpay create-order+verify+webhook (V39). ShedLock V41. Rate-limit V42. Local PAYMENT_MOCK default true; prod false.", WARN)

    # 03
    keep(divider(prs, 3, "Problem statement", "Safety, work, health, money and trust — fragmented."))
    s = keep(content(prs, "03  PROBLEM STATEMENT", "What women actually face (2026)", "MARKET RESEARCH with sources — not Youthian stats"))
    table(s, ["Problem", "Evidence", "Source / date"], [
        ["Registered VAW still high", "4,41,534 cases (2024); rate 64.6/lakh", "NCRB via The Print / CEDA, 2026"],
        ["112 exists but is not a livelihood OS", "ERSS in 36 States/UTs; Nirbhaya ₹7,712.85 cr allocated to FY25", "PIB Mar 2025"],
        ["Urban women still under-participate", "Urban female LFPR 24.8% CWS June 2026 vs rural 36.6%", "MoSPI PLFS bulletin June 2026"],
        ["Usual-status FLFPR 40% still << male 79.1%", "PLFS Annual Report CY 2025", "MoSPI Mar 2026 press note"],
        ["Gig/jobs apps don’t own SOS graph", "Sitha/Kaabil/Naukri have no police dispatch", "INFERENCE"],
        ["Trust / KYC theatre", "Many apps claim verified women; ops quality unpublished", "INFERENCE"],
        ["5–10 apps for one life", "112 + Naukri + Practo + Instagram + UPI + marketplace", "INFERENCE"],
    ], rh=0.52, fs=12)

    # 04
    keep(divider(prs, 4, "Objectives", "Measurable where possible — mix of CURRENT and RECOMMENDATION."))
    s = keep(content(prs, "04  OBJECTIVES", "What success looks like", "RECOMMENDATION — not current KPIs in code"))
    table(s, ["Objective", "Metric (proposed)", "Horizon"], [
        ["SOS reliability", "P95 notify < 30s; drill pass; never paywalled", "90 days"],
        ["Trust", "Admin SLA on verify/reject; audit log 100% of approvals", "90 days"],
        ["Safety + 112 coexistence", "In-app 112 deep-link + dual-app guidance", "6 months"],
        ["Livelihood density", "Verified supply in 1 metro (doctors/salons/workers/sellers)", "12 months"],
        ["Economic", "Take-rate on paid modules; SOS remains free", "Ongoing"],
        ["Inclusion", "Vernacular + DPDP deletion + accessibility", "6–12 months"],
        ["Honesty", "No AI marketing until LLM exists; dummy wallet removed", "90 days"],
    ], rh=0.55, fs=13)

    # 05
    keep(divider(prs, 5, "Project scope", "Current → required → future"))
    s = keep(content(prs, "05  PROJECT SCOPE", "Three horizons", "CURRENT IMPLEMENTATION vs RECOMMENDATION vs FUTURE"))
    card(s, 0.4, 1.22, 4.05, 5.7, "Current scope (ships today)",
         "Auth OTP/JWT. SOS+journey+contacts. Reels (web). Doctors, Glow, fitness, martial arts. Jobs/lawyers. Products+delivery. Events APIs. Funding portals. FL content. Creator hub. Admin queues. Razorpay path.", OK)
    card(s, 4.6, 1.22, 4.05, 5.7, "Development scope (required)",
         "Flyway completeness under validate. Real FCM send. SMS on. Audit log + RBAC split. Dummy wallet out. DPDP. 112 deep-link. Product returns on mobile. Panic/audio on Flutter. Secrets rotation. Tests/CI. Moderator people+tooling.", WARN)
    card(s, 8.8, 1.22, 4.1, 5.7, "Future scope",
         "Predictive safety (ethics). LLM copilot with consent. Intelligent matching. iOS. City expansion. CSR/gov. Safetipin-class maps partnership. Not: becoming a bank, AIF, or fake police PSAP.", MUTED)

    # 06
    keep(divider(prs, 6, "Target audience / users", "Actual UserType vs proposed roles."))
    s = keep(content(prs, "06  USERS", "Roles in code vs proposed", "FACT: enum has 13 values. Mentor/Employer/Moderator are NOT in UserType."))
    table(s, ["Group", "In code?", "How they appear"], [
        ["Woman / USER", "CURRENT", "UserType.USER"],
        ["ADMIN", "CURRENT", "Separate Admin entity + UserType.ADMIN"],
        ["VOLUNTEER", "CURRENT", "SOS respond"],
        ["DOCTOR / SALON / STYLIST / CENTRE / FITNESS_TRAINER", "CURRENT", "Portals + admin verify"],
        ["PROVIDER (incl. lawyer)", "CURRENT", "ProviderCategory.WOMEN_LAWYER etc."],
        ["SELLER / EVENT_HOST / ENTREPRENEUR / INVESTOR", "CURRENT", "Dedicated tables + APIs"],
        ["FinancialEducator / DeliveryPartner / Creator / Job worker", "CURRENT", "Tables/APIs — not UserType enum"],
        ["Mentor / Employer / Moderator / SuperAdmin / NGO", "FUTURE / proposed", "Do not document as live roles"],
    ], rh=0.5, fs=12)

    # 07 market
    keep(divider(prs, 7, "Market research 2026", "India + global + pricing + trends + needs. Sources cited."))
    s = keep(content(prs, "07  MARKET  ·  INDIA", "Official statistics", "MARKET RESEARCH"))
    table(s, ["Topic", "Figure", "Year", "Source"], [
        ["Crimes against women (registered)", "4,41,534; rate 64.6 / lakh", "2024", "NCRB / The Print / CEDA 2026"],
        ["Nirbhaya Fund allocation", "₹7,712.85 crore to FY2024–25", "PIB Mar 2025", "GoI PIB"],
        ["Nirbhaya utilisation", "~₹5,846 cr (~76%)", "PIB Mar 2025", "GoI PIB"],
        ["ERSS-112", "All 36 States/UTs", "PIB Mar 2025", "GoI"],
        ["WHL-181", "1.95 cr+ calls; 81.64 lakh+ women assisted", "PIB Mar 2025", "GoI — as stated"],
        ["Female LFPR usual status 15+", "40.0% (male 79.1%)", "CY 2025", "MoSPI PLFS Annual"],
        ["Female LFPR CWS 15+", "32.7% overall; urban 24.8%", "June 2026", "MoSPI monthly bulletin"],
        ["Lakhpati Didi", "2 cr+ achieved; target 3 cr by 2027", "Budget 2026 coverage", "The Hindu / MoRD as reported"],
        ["SHE-Marts", "Community SHG retail announced Budget 2026-27", "1 Feb 2026", "Union Budget / The Hindu"],
    ], rh=0.45, fs=11)

    s = keep(content(prs, "07  MARKET  ·  GLOBAL", "Safety-tech and family apps (listed where possible)", "Do not treat vendor TAM reports as audited."))
    table(s, ["Market", "Signal", "Implication for FightDFear"], [
        ["US / global family safety", "Life360 FY2025 revenue $489.5m; 95.8m MAU (31 Dec 2025). Q2’26: 102.4m MAU, qtr rev $159.0m. Source: Life360 IR / GlobeNewswire Mar & Aug 2026.", "Subscription + location graph works — but it is family, not women-economy."],
        ["US panic monitoring", "Noonlight: core SOS free US-only; Premium ~$9.99/mo reported (ImAlive Jul 2026).", "Paid monitoring is a US PSAP product. India equivalent is 112, not a startup PSAP."],
        ["UK", "Hollie Guard Extra £7.99/mo (ImAlive Jul 2026).", "URN/police partnership — learn trust, don’t copy pricing into INR blindly."],
        ["SEA / ME", "Women-only social/safety apps exist; scale usually Not publicly disclosed.", "Partnership > copycat."],
        ["Europe GDPR / India DPDP", "Location + SOS + ID docs are high-risk processing.", "Consent, retention, erasure are table stakes."],
    ], rh=0.72, fs=11)

    s = keep(content(prs, "07  MARKET  ·  PRICING", "Competitor pricing (public)", "Not publicly disclosed where blank"))
    table(s, ["Platform", "Model", "Public price", "Source"], [
        ["112 India / Himmat+ / 181", "Public good", "Free", "MHA / state police"],
        ["Life360 US", "Freemium sub", "Silver $7.99 / Gold $14.99 / Platinum $24.99 /mo (US pages 2026; some geos $9.99 Silver)", "life360.com"],
        ["Noonlight", "Free + premium", "Premium ~$9.99/mo reported", "ImAlive Jul 2026"],
        ["Hollie Guard Extra", "Subscription", "£7.99/mo", "ImAlive Jul 2026"],
        ["Sitha", "Marketplace + boosters", "Commission Not publicly disclosed", "Company site / Hindu BL"],
        ["Info Edge (Naukri parent)", "B2B recruitment", "Employer pricing Not a consumer SKU", "FY26 IR"],
        ["FightDFear CURRENT", "Take-rate fields in schema", "No complete consumer sub SKU", "Code"],
    ], rh=0.52, fs=11)

    s = keep(content(prs, "07  MARKET  ·  TRENDS & NEEDS", "What women are looking for", "INFERENCE grounded in official + competitor evidence"))
    card(s, 0.4, 1.22, 4.05, 2.5, "Safety tech", "Official 112 + private SOS + journey share. Wearables/family circles (Life360) prove location graphs monetize — ethically SOS must stay free in India.", ROSE)
    card(s, 4.6, 1.22, 4.05, 2.5, "Work & money", "Urban FLFPR still low. Sitha/SheJobs/Kaabil/SHE-Marts show political + commercial push for women gigs and SHG retail.", TEAL)
    card(s, 8.8, 1.22, 4.1, 2.5, "Trust & AI", "Info Edge FY26: AI in matching (MD commentary). Sitha claims AI matching. FightDFear has no LLM — do not fake it. Verification ops beat AI theatre.", NAVY)
    card(s, 0.4, 3.9, 12.5, 2.95, "User needs (ranked for FightDFear)",
         "1) A panic path that actually reaches someone  2) Verified local doctors/salons/workers  3) Flexible income  4) Privacy of location and ID  5) Community that is moderated  6) Convenience of one identity  7) Personalization later. Convenience without trust is a liability.", WARN)

    # 08 existing
    keep(divider(prs, 8, "Existing features", "Only what the repository implements."))
    s = keep(content(prs, "08  EXISTING FEATURES  ·  1/3", "Safety, identity, community", "CURRENT IMPLEMENTATION"))
    table(s, ["Feature", "Purpose", "Web", "Mobile", "Admin"], [
        ["Register / login / OTP / JWT", "Identity", "Session JSP", "FULLY", "Users verify/ban"],
        ["SOS trigger/history", "Emergency", "FULLY", "FULLY", "List/resolve"],
        ["Trusted / emergency contacts", "Notify graph", "FULLY", "FULLY", "—"],
        ["Journey + ShedLock escalate", "Travel safety", "FULLY", "FULLY", "—"],
        ["Buddy", "Companion", "PARTIAL", "PARTIAL", "Buddy mgmt"],
        ["Panic / recording / battery / SOS audio", "Evidence", "JSP/API", "MISSING", "—"],
        ["Danger points / heatmap", "Map intel", "PARTIAL", "Danger map", "Verify points"],
        ["Reels / comments / follow", "Community", "FULLY", "PARTIAL", "Block reported"],
        ["Chat", "Messaging", "FULLY", "PARTIAL", "—"],
        ["Groups / express posts", "Social", "JSP", "MISSING", "—"],
    ], fs=11)

    s = keep(content(prs, "08  EXISTING FEATURES  ·  2/3", "Livelihood & care", "CURRENT IMPLEMENTATION"))
    table(s, ["Feature", "Web", "Mobile", "Backend/API", "Admin"], [
        ["Doctors + Razorpay + chat", "FULLY", "FULLY", "/api/doctors*", "Verify/reject/changes"],
        ["Glow salons/stylists/bookings", "FULLY", "FULLY", "/api/glow*", "Approve salon/stylist"],
        ["Fitness trainers", "FULLY", "FULLY", "/api/fitness*", "Verify/suspend"],
        ["Martial arts + certificates", "FULLY", "FULLY", "/api/martial-arts*", "Approve centres"],
        ["Jobs / worker bookings", "FULLY", "FULLY", "/api/marketplace/jobs", "Approve applications"],
        ["Lawyers / ServiceProvider", "FULLY", "FULLY", "/api/marketplace", "Pending providers"],
        ["Products cart/order/track", "FULLY", "FULLY minus returns", "/api/women-products", "Sellers + orders"],
        ["Delivery partners", "PARTIAL", "FULLY portal", "/api/delivery", "Approve partners"],
    ], fs=11)

    s = keep(content(prs, "08  EXISTING FEATURES  ·  3/3", "Growth modules, money, AI, push", "CURRENT IMPLEMENTATION"))
    table(s, ["Feature", "Status", "Notes"], [
        ["Events + hosts", "PARTIAL", "APIs+Flutter+admin; prod table missing until V45"],
        ["Entrepreneur / investor / proposals", "PARTIAL", "Portals + admin approve + investment release"],
        ["Financial literacy", "PARTIAL", "Videos/sessions/workshops/loans JSP — not regulated advice"],
        ["Creator hub tips/cashout", "FULLY core", "Admin cashout; T&S still ops"],
        ["Wallet redeem", "PLACEHOLDER", "Hardcoded ‘10% Off Salon…’ + fake COUPON-timestamp"],
        ["Razorpay", "PARTIAL", "Real verify path; local mock default true"],
        ["Email", "FULLY", "OTP, reset, SOS async — secrets committed in properties"],
        ["SMS", "PARTIAL", "Twilio/Msg91; sms.enabled=false default"],
        ["FCM", "PLACEHOLDER", "Stores tokens, logs payload — HTTP v1 not wired"],
        ["AISafetyService Weka", "PLACEHOLDER", "Needs crime_data.arff / safety_model.model"],
        ["OpenAI / Gemini", "MISSING", "No matches in Java or Dart"],
    ], rh=0.4, fs=12)

    # 09 drawbacks
    keep(divider(prs, 9, "Existing system drawbacks", "Honest, specific, from code."))
    s = keep(content(prs, "09  DRAWBACKS", "What is actually wrong today", "FACT from SecurityConfig, services, Flyway, wallets"))
    table(s, ["Gap class", "Specific issue", "Priority"], [
        ["Security", "Committed mail/Razorpay/Maps/JWT defaults; JwtUtil hardcoded fallback; Maps key in Flutter", "P0"],
        ["Security", "CSRF disabled; huge PUBLIC_URLS; no hasRole; /admin/registerAdmin public pattern", "P0"],
        ["Schema/ops", "Prod ddl-auto=validate vs Hibernate-era tables; Flyway repair() every startup", "P0"],
        ["SOS", "FCM not sent; SMS off; no 112; panic/audio not on Flutter", "P0"],
        ["Admin", "God-admin; no audit log entity; no Moderator", "P0"],
        ["Money", "Coins + walletBalance + payout_balance + dummy coupons", "P1"],
        ["Mobile parity", "Groups, returns, recording, battery missing on Flutter", "P1"],
        ["Quality", "Few automated tests vs ~108 controllers / ~175 entity types", "P0"],
        ["Legal", "Health + funding + FL without strong in-product disclaimers", "P0"],
        ["Analytics", "Dashboard.meta exists; not a product analytics stack", "P2"],
    ], rh=0.45, fs=12)

    # 10 unique
    keep(divider(prs, 10, "Unique features", "Ordinary features are not unique."))
    s = keep(content(prs, "10  DIFFERENTIATION", "Existing vs potential", "Do not call SOS unique — 112 already does dispatch."))
    card(s, 0.4, 1.22, 6.2, 5.7, "Existing differentiation (weak but real)",
         "One codebase already spans SOS graph + verified doctors/salons/jobs/products/events/funding/creator. Admin verification factory is broader than typical SOS apps. Unified User identity across verticals exists in schema. That combination is rare in India — execution quality is not yet a moat.", ROSE)
    card(s, 6.8, 1.22, 6.1, 5.7, "Potential differentiation (not live)",
         "Safety-to-career journey. Reputation that travels across modules. Cross-module rewards (today dummy). Women-business network with KYC that actually operates. AI assistance (MISSING). 112 coexistence UX. City-level density. CSR/Nirbhaya-adjacent partnerships without replacing the state.", NAVY)

    # 11 must have
    keep(divider(prs, 11, "Features to include", "MoSCoW for a competitive product."))
    s = keep(content(prs, "11  FEATURES TO INCLUDE", "MoSCoW", "RECOMMENDATION"))
    table(s, ["Priority", "Items"], [
        ["MUST", "SOS drill + real push/SMS · 112 deep-link · secrets/CSRF/RBAC/audit · Flyway green · DPDP deletion · dummy wallet gone · health/funding disclaimers · moderator queue"],
        ["SHOULD", "Mobile parity (panic/audio/returns/groups) · unified ledger · vernacular · forced-update · crash reporting · settlement reports"],
        ["COULD", "Partner boosters · event tickets · employer job slots · creator ads (safe) · Safetipin data partnership"],
        ["FUTURE", "LLM copilot · predictive maps · iOS · international · accredited learning"],
    ], rh=1.15, fs=13)

    # 12 future
    keep(divider(prs, 12, "Future enhancements", "Clearly FUTURE — not in the APK."))
    s = keep(content(prs, "12  FUTURE ENHANCEMENTS", "After P0 trust", "FUTURE"))
    bullets(s, [
        "Advanced AI: moderation assist, matching, safety copilot — only with consent; never send SOS to third-party LLMs by default.",
        "Predictive safety: partnership data (Safetipin-class), not rumor maps.",
        "Intelligent job matching / personalized learning / AI mentorship — after supply density.",
        "BI + safety intelligence dashboards for city partners.",
        "Government/NGO: 112, 181, OSC, SHE-Marts/Lakhpati Didi as referrals — FightDFear is not the scheme operator.",
        "International: only after one Indian metro is trustworthy.",
    ], size=15)

    # 13 competitors
    keep(divider(prs, 13, "Competitors analysis", "2026 landscape — not a copied Youthian list."))
    s = keep(content(prs, "13  COMPETITORS  ·  SAFETY", "Direct and adjacent", "MARKET RESEARCH"))
    table(s, ["Platform", "Country", "Purpose", "Model", "Learning"], [
        ["112 India / ERSS", "India", "Official dispatch", "Public / Nirbhaya", "Do not compete; deep-link"],
        ["Himmat Plus", "India (Delhi+)", "Police SOS", "Free", "Trust brand, geo-limited"],
        ["Safetipin", "India/global", "Safety scores", "Freemium / B2G", "Maps intel partnership"],
        ["Life360", "US/global", "Family circles", "Subscription", "Location graph + paid tiers"],
        ["Noonlight / bSafe / Hollie", "US/UK", "Panic + monitor", "Free+sub", "PSAP is not India’s gap"],
        ["Raksha / VithU / Shake2Safety", "India", "Simple SOS", "Free/ad", "Churn without livelihood"],
    ], rh=0.55, fs=12)

    s = keep(content(prs, "13  COMPETITORS  ·  ECONOMY", "Jobs, gigs, commerce, careers", "MARKET RESEARCH"))
    table(s, ["Platform", "Purpose", "Users / $", "Weakness vs FDF"], [
        ["Sitha", "Women gig + products", "Goal 1M entrepreneurs by 2027 (company, Hindu BL). SheJobs ~1 lakh DB (company-stated).", "No SOS graph"],
        ["SHE-Marts (Budget 2026)", "SHG community retail", "Policy, not an app competitor yet", "Physical retail; different layer"],
        ["Info Edge / Naukri", "Recruitment", "Standalone FY26 rev ₹3,052 cr (company PR 22 May 2026)", "Not women-safety; B2B"],
        ["Kaabil / HerJobs", "Women jobs", "Kaabil scale Not publicly disclosed as MAU", "No emergency"],
        ["Practo / 1mg", "Health", "Not publicly disclosed here", "Not women-only OS"],
        ["Instagram / YT", "Attention", "Meta — not comparable MAU mix", "No KYC livelihood or SOS"],
    ], rh=0.58, fs=11)

    # 14 comparison
    keep(divider(prs, 14, "Feature comparison", "FightDFear column = code status."))
    s = keep(content(prs, "14  FEATURE COMPARISON", "Safety + opportunity matrix", "CURRENT IMPLEMENTATION vs public product positioning"))
    table(s, ["Feature", "FightDFear", "112", "Life360", "Sitha", "Naukri"], [
        ["SOS / dispatch", "Private contacts PARTIAL", "Official", "Circle + assist", "—", "—"],
        ["Journey / live location", "PARTIAL", "Trip/SOS", "Strong", "Job track", "—"],
        ["Community video", "Web FULLY", "—", "—", "—", "—"],
        ["Jobs / gigs", "FULLY core", "—", "—", "Strong", "Strong"],
        ["Doctors / beauty", "FULLY core", "—", "—", "Wellness gigs", "—"],
        ["Products", "PARTIAL", "—", "—", "Handmade", "—"],
        ["Funding / FL", "PARTIAL", "—", "—", "—", "—"],
        ["AI / LLM", "MISSING", "—", "Some", "Matching claims", "AI matching (FY26 commentary)"],
        ["Admin verification", "FULLY queues", "Govt ID", "Family", "KYC claimed", "Employer"],
        ["Wallet / rewards", "PLACEHOLDER", "—", "Sub", "Payments", "—"],
    ], fs=11)

    # 15 pricing
    keep(divider(prs, 15, "Pricing", "Theirs vs proposed ours. SOS never paywalled."))
    s = keep(content(prs, "15  PRICING", "Proposed FightDFear commercial offers", "RECOMMENDATION — not live SKUs"))
    table(s, ["Offer", "Who pays", "Draft", "Rule"], [
        ["Free core", "User", "SOS, contacts, journey, 112 link, basic feed", "Never meter panic"],
        ["Premium membership", "User", "Later: extra storage, boosts, FL certificates — not SOS", "After trust"],
        ["Provider plan / booster", "Salon/doctor/worker/seller", "Featured listing + analytics", "Like Sitha boosters"],
        ["Marketplace commission", "Seller/provider", "Use existing commission_percent / platform_fee fields", "Transparent"],
        ["Employer / job slots", "Business", "Paid job posts — Employer role is FUTURE", "Don’t spam workers"],
        ["Events", "Host/attendee", "Ticket take-rate after schema stable", "P1"],
        ["Enterprise / CSR / city", "NGO/Gov/Corp", "Safe-city dashboard, verified supply", "Procurement-long"],
    ], rh=0.55, fs=12)

    # 16 revenue
    keep(divider(prs, 16, "Revenue of existing websites", "Verified vs estimated. Never fabricated."))
    s = keep(content(prs, "16  COMPETITOR REVENUE", "Public filings only", "Verified = company IR / statutory. Else Not publicly disclosed."))
    table(s, ["Company", "Revenue", "Year", "Source", "Model", "Type"], [
        ["Life360", "US$489.5 million", "FY 2025", "Life360 Q4’25 / GlobeNewswire 2 Mar 2026", "Sub + hardware + other", "Verified"],
        ["Life360", "US$159.0m quarter; H1 US$302.1m", "Q2 / H1 2026", "Life360 IR 10 Aug 2026", "Same", "Verified"],
        ["Info Edge (standalone)", "₹3,052 crore", "FY2025-26", "Company PR 22 May 2026", "Recruitment + 99acres + JS + education", "Verified"],
        ["Info Edge recruitment", "₹2,256 crore segment (results PDF)", "FY26", "Audited results 31 Mar 2026", "B2B jobs", "Verified"],
        ["Sitha / SheJobs / Kaabil / Safetipin / Practo", "—", "—", "—", "—", "Not publicly disclosed"],
        ["112 India", "N/A (public service)", "—", "Nirbhaya funded", "Taxpayer", "N/A"],
        ["FightDFear", "—", "2026", "Private", "Take-rate nascent", "Not publicly disclosed"],
    ], rh=0.5, fs=11)

    # 17 workflow
    keep(divider(prs, 17, "Workflow", "As coded — not aspirational ATS."))
    s = keep(content(prs, "17  WORKFLOWS  ·  1/2", "Identity and safety", "CURRENT IMPLEMENTATION"))
    tb(s, Inches(0.45), Inches(1.2), Inches(12), Inches(0.3), "Registration (mobile)", 14, True, NAVY)
    flow_row(s, ["OTP / register", "JWT", "Landing feed", "Dashboard", "Profile"], 1.55)
    tb(s, Inches(0.45), Inches(2.45), Inches(12), Inches(0.3), "Partner", 14, True, NAVY)
    flow_row(s, ["register-quick", "Email OTP", "Profile complete", "Admin verify", "Dashboard"], 2.8)
    tb(s, Inches(0.45), Inches(3.7), Inches(12), Inches(0.3), "SOS  (no police PSAP)", 14, True, ROSE)
    flow_row(s, ["Trigger", "GPS payload", "Email contacts", "SMS if on", "Admin resolve"], 4.05)
    tb(s, Inches(0.45), Inches(5.3), Inches(12.4), Inches(1.5),
       "GAP: FCM is log-only. SMS default off. Panic timer is JSP. 112 is MISSING. RECOMMENDATION: dual-path 112 + private graph.",
       13, False, MUTED)

    s = keep(content(prs, "17  WORKFLOWS  ·  2/2", "Economy modules", "CURRENT IMPLEMENTATION"))
    tb(s, Inches(0.45), Inches(1.2), Inches(12), Inches(0.28), "Healthcare / Glow / Fitness", 13, True, NAVY)
    flow_row(s, ["Verify partner", "Availability", "Book", "Razorpay", "Chat / consult"], 1.52)
    tb(s, Inches(0.45), Inches(2.45), Inches(12), Inches(0.28), "Jobs  (no Employer UserType)", 13, True, NAVY)
    flow_row(s, ["Worker OTP", "Apply", "Admin may approve", "Booking", "Review"], 2.78)
    tb(s, Inches(0.45), Inches(3.7), Inches(12), Inches(0.28), "Marketplace products", 13, True, NAVY)
    flow_row(s, ["Seller verify", "Catalog", "Order", "Delivery track", "Return (web)"], 4.02)
    tb(s, Inches(0.45), Inches(5.25), Inches(12), Inches(0.28), "Funding  (not an AIF)", 13, True, NAVY)
    flow_row(s, ["Proposal", "Admin approve", "Investor", "Chat / meet", "Release $"], 5.55)

    # 18 user flow
    keep(divider(prs, 18, "User flow", "Journeys for each persona that exists."))
    s = keep(content(prs, "18  USER JOURNEYS", "Personas mapped to code", "Mentor/Employer journeys are FUTURE"))
    table(s, ["Persona", "Journey (as coded)", "Gap"], [
        ["Normal user", "Auth → landing → discover modules → engage reels/book/buy", "Feed ranking thin"],
        ["Emergency user", "SOS → location in payload → contacts → admin resolve", "No 112, weak push"],
        ["Job seeker", "Jobs OTP → apply → bookings", "No interview ATS"],
        ["Employer", "FUTURE role", "Use worker/admin today"],
        ["Seller", "Seller auth → products → orders → delivery", "Mobile returns MISSING"],
        ["Mentor", "FUTURE", "Not a UserType"],
        ["Healthcare provider", "Doctor OTP → complete profile → admin → appointments", "Clinical/legal SOP"],
        ["Entrepreneur", "Register → proposal → admin → investor", "Securities disclaimer"],
    ], rh=0.52, fs=12)

    # 19 architecture
    keep(divider(prs, 19, "System architecture", "Current vs recommended production."))
    s = keep(content(prs, "19  ARCHITECTURE  ·  CURRENT", "Monolith as built", "FACT"))
    labels = ["Flutter JWT\nJSP session\nAdmin", "Spring Boot\n3.4 :8084", "Modules\n108 controllers", "MySQL 8\nFlyway V45", "Razorpay\nSMTP Maps"]
    xs = [0.4, 3.0, 5.6, 8.2, 10.8]
    for i, lab in enumerate(labels):
        rect(s, Inches(xs[i]), Inches(2.0), Inches(2.15), Inches(1.7), NAVY if i % 2 == 0 else ROSE, radius=True)
        tb(s, Inches(xs[i]), Inches(2.35), Inches(2.15), Inches(1.2), lab, 13, True, WHITE, PP_ALIGN.CENTER)
        if i < 4:
            tb(s, Inches(xs[i] + 2.05), Inches(2.5), Inches(0.3), Inches(0.5), ">", 18, True, NAVY, PP_ALIGN.CENTER)
    tb(s, Inches(0.45), Inches(4.0), Inches(12.4), Inches(2.8),
       "Auth: JWT filter + session hydrate. CSRF off. PUBLIC_URLS permitAll then authenticated().\n"
       "Cache/search: none as products (DB queries). Storage: local uploads (see STORAGE_MIGRATION.md).\n"
       "Realtime: STOMP/WebSocket for some chats. Push: not actually FCM.\n"
       "Weakness: two auth modes, schema validate incidents, Flyway repair on boot, committed secrets.",
       14, False, MUTED)

    s = keep(content(prs, "19  ARCHITECTURE  ·  RECOMMENDED", "Production target", "RECOMMENDATION / FUTURE"))
    bullets(s, [
        "Keep monolith until SOS+payments+schema are boringly stable — do not rewrite for fashion.",
        "Secrets only via env; delete committed keys; rotate JWT; Maps via backend proxy.",
        "Spring method security + Moderator/Finance/SuperAdmin; CSRF for cookie session or go JWT-only for APIs.",
        "Stop Flyway repair-on-start. ddl-auto=validate forever in prod. CREATE TABLE migrations for remaining Hibernate-era tables.",
        "Object storage ACL for ID docs. Unified ledger. Real FCM HTTP v1. SMS failover.",
        "OpenAPI /api/v1. APM + SOS SLO dashboards. Optional Redis later — not a current requirement to ‘look cloud native’.",
        "112 deep-link + documented dual-app. No fake PSAP.",
    ], size=15)

    # 20 modules
    keep(divider(prs, 20, "Modules overview", "Purpose, status, gaps."))
    s = keep(content(prs, "20  MODULE MAP", "Complete inventory", "CURRENT IMPLEMENTATION"))
    table(s, ["Module", "Users", "Status", "Primary API / UI", "Gap"], [
        ["SOS / journey / buddy", "USER, VOLUNTEER, ADMIN", "PARTIAL", "/api/sos /journey /buddy", "112, FCM, Flutter panic"],
        ["Community", "USER", "PARTIAL", "/video /api/videos /chat", "Mobile groups, mods"],
        ["Doctors", "USER, DOCTOR", "FULLY", "/api/doctors", "Clinical SOP"],
        ["Glow", "USER, SALON, STYLIST", "FULLY", "/api/glow", "Settlement"],
        ["Fitness / martial", "USER, trainer, CENTRE", "FULLY", "/api/fitness /martial-arts", "—"],
        ["Jobs / lawyers", "Worker, PROVIDER", "FULLY", "/api/marketplace", "No Employer"],
        ["Products / delivery", "USER, SELLER", "PARTIAL", "/api/women-products /delivery", "Mobile returns"],
        ["Events", "USER, EVENT_HOST", "PARTIAL", "/api/women-events", "Schema ops"],
        ["Funding", "ENTREPRENEUR, INVESTOR", "PARTIAL", "/api/entrepreneur /investor", "Legal"],
        ["FL / creator / wallet", "Educator, creator, USER", "PARTIAL", "/api/financial-literacy /creator-hub /wallet", "Dummy coins, advice line"],
        ["Admin", "ADMIN", "FULLY queues", "/admin", "RBAC, audit"],
    ], fs=10)

    # 21 RBAC
    keep(divider(prs, 21, "User roles & permissions", "Actual matrix — ADMIN is omnipotent."))
    s = keep(content(prs, "21  RBAC", "What code allows today", "FACT: Spring does not hasRole on UserType"))
    table(s, ["Action", "USER", "Partner*", "VOLUNTEER", "ADMIN"], [
        ["View own data", "Y", "Y", "Y", "Y"],
        ["Create SOS / book / apply / order", "Y", "supply side", "respond", "Y"],
        ["Verify / reject / request-changes", "—", "—", "—", "Y"],
        ["Ban / suspend", "—", "—", "—", "Y"],
        ["Moderate reels", "report", "creator extras", "—", "block"],
        ["Release investment / cashout", "—", "request", "—", "Y"],
        ["Configure platform", "—", "—", "—", "Y (unsplit)"],
        ["Moderator / Finance only", "MISSING", "MISSING", "MISSING", "MISSING"],
    ], rh=0.48, fs=12)
    tb(s, Inches(0.4), Inches(6.55), Inches(12.4), Inches(0.45),
       "*Partner = DOCTOR, SALON/STYLIST, CENTRE, PROVIDER, SELLER, EVENT_HOST, FITNESS_TRAINER, ENTREPRENEUR, INVESTOR + educator/delivery/creator tables.",
       11, False, MUTED)

    # 22 tech
    keep(divider(prs, 22, "Technology stack", "Currently used vs recommended."))
    s = keep(content(prs, "22  TECHNOLOGY", "Extracted from the repo", "Never invented"))
    table(s, ["Layer", "Currently used", "Recommended"], [
        ["Mobile", "Flutter, Provider, geolocator, google_maps_flutter, razorpay_flutter", "Crashlytics/Sentry, pinning, iOS later"],
        ["Web", "JSP + some static HTML; Thymeleaf unused", "Don’t migrate templates mid-crisis"],
        ["Backend", "Spring Boot 3.4.3, Java 17, WAR", "Method security, OpenAPI"],
        ["DB", "MySQL 8, JPA, Flyway V45", "No ddl-auto=update in any shared env"],
        ["Auth", "BCrypt, JWT, session, email OTP", "Rotate secrets; CSRF or JWT-only APIs"],
        ["Pay", "Razorpay + V39 fulfillment", "Keep mock false in prod"],
        ["Jobs", "ShedLock, rate_limit_buckets", "SOS SLO monitors"],
        ["AI", "Weka J48 placeholder", "LLM only after DPDP + human-in-loop"],
        ["Push/SMS", "Token store; SMS flag off", "FCM HTTP v1 + Msg91/Twilio on"],
        ["Maps", "Google Maps key in app + server", "Restrict keys; proxy"],
    ], rh=0.42, fs=11)

    # 23 similar apps
    keep(divider(prs, 23, "Similar apps — India & global", "Why each is relevant."))
    s = keep(content(prs, "23  SIMILAR APPS", "Relevance map", "MARKET RESEARCH"))
    table(s, ["Region", "Platform", "Why relevant"], [
        ["India", "112 India, Himmat Plus, Safetipin", "Official/civic safety layer FDF must coexist with"],
        ["India", "Sitha, SheJobs, Kaabil, Naukri", "Livelihood competitors / distribution"],
        ["India", "SHE-Marts (policy 2026), Lakhpati Didi", "State SHG commerce — partner, don’t clone poorly"],
        ["India", "Practo, Urban Company (indirect)", "Health and home services UX bar"],
        ["Global", "Life360, Noonlight, bSafe, Hollie Guard", "SOS product polish and paid monitoring"],
        ["Global", "Instagram, LinkedIn (indirect)", "Community and career attention markets"],
    ], rh=0.62, fs=13)

    # 24 comparison 2
    keep(divider(prs, 24, "Comparison with features", "Second matrix — trust layer."))
    s = keep(content(prs, "24  TRUST & VERTICALS", "Where FightDFear can win or lose", "CURRENT vs public positioning"))
    table(s, ["Layer", "FDF", "Specialists win if…", "FDF wins if…"], [
        ["Safety", "PARTIAL private SOS", "User needs police now → 112", "Private graph + journey + 112 together"],
        ["Community", "PARTIAL", "Instagram attention", "Moderated women-only + reputation"],
        ["Jobs", "FULLY core", "Naukri/Kaabil scale", "Verified + safer context"],
        ["Care (doc/glow/fit)", "FULLY core", "Practo/UC density", "Women-verified local"],
        ["Commerce", "PARTIAL", "Amazon/SHE-Marts", "Women-owned + KYC ops"],
        ["Funding/FL", "PARTIAL", "Banks/SEBI platforms", "Literacy + discovery only"],
        ["Admin/T&S", "FULLY queues / MISSING people", "Anyone with a Trust org", "Queues + staffed SLAs"],
        ["AI", "MISSING", "Anyone shipping matching", "Don’t fake; add later"],
    ], rh=0.5, fs=12)

    # 25 users
    keep(divider(prs, 25, "How many users they have", "Metric type matters."))
    s = keep(content(prs, "25  USER / SCALE METRICS", "Do not confuse downloads with MAU", "MARKET RESEARCH"))
    table(s, ["Platform", "Number", "Metric type", "When", "Source"], [
        ["Life360", "95.8 million", "MAU", "31 Dec 2025", "Life360 Q4’25"],
        ["Life360", "102.4 million", "MAU", "30 Jun 2026", "Life360 Q2’26 IR"],
        ["Life360", "2.8 million", "Paying Circles", "31 Dec 2025", "Life360 Q4’25"],
        ["112 India app", "~5.6 million", "Downloads (not MAU)", "Aug 2026", "AppBrain — third-party, not MHA"],
        ["SheJobs", "~1 lakh", "Database (company-stated)", "2025 Hindu BL", "Company-stated"],
        ["Sitha", "1 million by 2027", "Goal, not current users", "2025", "Company / Hindu BL"],
        ["Lakhpati Didi", "2 crore+", "Scheme members (not an app MAU)", "2026 coverage", "MoRD via The Hindu"],
        ["Safetipin / Kaabil / Himmat+ / FDF", "—", "—", "—", "Not publicly disclosed"],
    ], rh=0.48, fs=11)

    # 26 popular
    keep(divider(prs, 26, "How they got popular", "Then: what FightDFear can learn."))
    s = keep(content(prs, "26  GROWTH MECHANICS", "Evidence-based patterns", "INFERENCE on learnings; facts on listed firms"))
    table(s, ["Player", "How they scaled", "FightDFear learning"], [
        ["112 / ERSS", "Nirbhaya funding, state ERC mandate, free, power-button SOS", "Coexist; never pick a fight with the state PSAP"],
        ["Life360", "Family network effects, freemium, hardware, ads (2026 IR), 180+ countries", "Graph + paid tiers work; India women-economy is a different graph"],
        ["Naukri / Info Edge", "Two-sided B2B, 25+ years, cash machine (FY26 ₹1,469 cr CFO from ops standalone)", "Don’t attack white-collar recruitment head-on"],
        ["Sitha / SheJobs", "Founder narrative, state minister launch, women-only brand", "City campaigns + verification story"],
        ["Instagram", "Creators, network effects", "Reels without T&S is a liability"],
        ["SHG / Lakhpati / SHE-Marts", "Government program distribution", "CSR/govt as channel, not as product identity"],
    ], rh=0.58, fs=11)

    # 27 mistakes
    keep(divider(prs, 27, "What mistakes they did", "Verified issue vs strategic inference."))
    s = keep(content(prs, "27  COMPETITOR MISTAKES", "Separate evidence from inference", "Do not smear without source"))
    table(s, ["Observation", "Type", "Detail"], [
        ["112 India Play reviews: login/OTP failures (2025–26 comments)", "Verified issue (anecdotal reviews)", "Reliability of official app is a user complaint class — FDF still must deep-link, not mock"],
        ["Safety-only consumer apps stall on retention", "Strategic inference", "No income loop"],
        ["Women-gig apps over-claim ‘AI matching’", "Strategic inference", "Matching quality unpublished"],
        ["Family trackers (Life360) raise privacy debates globally", "Strategic inference + long-running public debate", "Location retention must be explicit under DPDP"],
        ["Marketplace KYC theatre", "Strategic inference", "Admin queues without people = fake trust"],
        ["Super-app complexity", "Strategic inference", "FDF already has this risk internally"],
    ], rh=0.62, fs=12)

    # 28 how we can
    keep(divider(prs, 28, "How we can do", "What FightDFear should do differently."))
    s = keep(content(prs, "28  HOW WE WIN DIFFERENTLY", "Playbook", "RECOMMENDATION"))
    card(s, 0.4, 1.22, 4.05, 2.55, "Product", "SOS free + 112. One verified identity. Kill dummy wallet. City density over 12 half-built verticals.", ROSE)
    card(s, 4.6, 1.22, 4.05, 2.55, "Trust", "Staffed admin SLAs, audit logs, CSAM process, DPDP. Verification is an operations product.", NAVY)
    card(s, 8.8, 1.22, 4.1, 2.55, "Tech", "Stabilize schema/secrets/push before AI. Honesty in status labels.", TEAL)
    card(s, 0.4, 3.95, 4.05, 2.85, "Go-to-market", "One metro. Campus + SHG/CSR + doctor/salon supply. Not national Super App ads.", WARN)
    card(s, 4.6, 3.95, 4.05, 2.85, "Monetize", "Take-rate and boosters. CSR/city. Never panic paywall.", ROSE)
    card(s, 8.8, 3.95, 4.1, 2.85, "Do not", "Fight Instagram, Naukri, UPI, or 112. Do not advertise LLM. Do not become a fund.", NAVY2)

    # 29 lacking
    keep(divider(prs, 29, "What people are looking for", "Unmet needs."))
    s = keep(content(prs, "29  UNMET NEEDS", "Largest opportunities", "MARKET RESEARCH + INFERENCE"))
    table(s, ["Need", "What exists", "What lacks", "Opportunity size for FDF"], [
        ["Immediate help", "112, 181, OSC", "Private graph + journey in same app as livelihood", "P0 coexist"],
        ["Flexible work", "Sitha, UC, Naukri, Kaabil", "Women-verified + safer context", "P1 city supply"],
        ["Trusted local care", "Practo, clinics", "Women-doctor/salon KYC ops", "P1 (already coded)"],
        ["Sell / SHG", "Amazon, SHE-Marts policy", "Digital storefront with real KYC", "P2"],
        ["Community without abuse", "IG", "Moderated women space", "P0 T&S"],
        ["One identity", "Nothing complete", "Fragmentation", "Core bet"],
        ["Honest AI", "Claims everywhere", "Few disclose limits", "Differentiate by not lying"],
    ], rh=0.52, fs=12)

    # 30 solution
    keep(divider(prs, 30, "How we can get a solution", "Problem → gap → FightDFear move."))
    s = keep(content(prs, "30  SOLUTION MAP", "Safety + trust + opportunity", "RECOMMENDATION"))
    table(s, ["Problem", "Existing gap", "FightDFear solution", "Priority"], [
        ["Fear in transit", "SOS apps don’t pay bills; 112 isn’t a companion", "Journey + SOS + 112 link", "P0"],
        ["Untrusted providers", "KYC theatre", "Admin queues + audit + SLA", "P0"],
        ["Urban under-employment", "FLFPR urban 24.8% CWS Jun 2026", "Jobs/Glow/products density in 1 city", "P1"],
        ["Fragmented apps", "5–10 logins", "Unified identity (already in schema)", "P1"],
        ["Fake AI", "Marketing vs code", "No LLM claim until built", "P0 brand"],
        ["Health/funding legal", "Blurred advice/securities", "Disclaimers; literacy not origination", "P0"],
        ["Retention", "Safety-only churn", "Earn + heal + belong after SOS works", "P1"],
    ], rh=0.52, fs=11)

    # 31 execution
    keep(divider(prs, 31, "Project execution plan", "Eight phases."))
    s = keep(content(prs, "31  EXECUTION", "Phased plan", "RECOMMENDATION — acceptance = tests + drills, not slideware"))
    table(s, ["Phase", "Focus", "Priority", "Accept"], [
        ["1 Audit & stabilize", "Secrets, CSRF plan, Flyway green, tests, mock=false", "P0", "validate boot; no committed secrets"],
        ["2 Core safety", "FCM+SMS, SOS drill, 112 link, Flutter panic/audio", "P0", "Drill runbook signed"],
        ["3 Community T&S", "Moderator role, report queue, CSAM policy", "P0", "SLA + audit log"],
        ["4 Jobs & career", "Employer story, dispute, quality supply", "P1", "City fill rate"],
        ["5 Education / healthcare", "SOP, disclaimers, licence ops", "P1", "Legal review"],
        ["6 Marketplace / entrepreneurship", "Ledger, returns mobile, funding disclaimer", "P1", "Settlement report"],
        ["7 AI / intelligence", "LLM + matching after privacy", "P3", "Consent + evals"],
        ["8 Scale & commercialize", "Boosters, CSR, 2nd city", "P2", "Unit economics"],
    ], rh=0.48, fs=11)

    # 32 team
    keep(divider(prs, 32, "Team needed", "MVP vs growth vs scale."))
    s = keep(content(prs, "32  TEAM", "Do not hire an ML org before SOS works", "RECOMMENDATION"))
    table(s, ["Role", "MVP", "Growth", "Scale"], [
        ["PM + BA", "1 combined", "1+1", "2 PM"],
        ["UI/UX", "1", "1–2", "2"],
        ["Flutter", "2", "3", "4"],
        ["Backend (Spring)", "2", "3", "5"],
        ["DB / DevOps", "1", "1–2", "2 + cloud"],
        ["QA", "1", "2", "3"],
        ["Security", "part-time", "1", "1 + pen-test retainer"],
        ["Trust & Safety / support", "1 (critical)", "3", "shift coverage"],
        ["AI/ML / data", "0", "0–1", "2"],
        ["Growth / community / marketing", "0.5", "2", "4"],
    ], rh=0.42, fs=12)

    # 33 commercialize
    keep(divider(prs, 33, "Commercialization", "How the product makes money without endangering users."))
    s = keep(content(prs, "33  COMMERCIALIZATION  ·  1/2", "Revenue architecture", "RECOMMENDATION"))
    table(s, ["Motion", "Stream", "When"], [
        ["B2C", "Free SOS forever. Premium later (storage/boosts) — not emergency", "After trust"],
        ["B2C education", "FL certificates / workshop fees (content, not advice)", "P2"],
        ["B2B recruitment", "Paid job slots when Employer exists", "P2"],
        ["B2B provider", "Salon/doctor/trainer subscriptions + featured", "P1"],
        ["Marketplace", "Commission (schema exists) + seller plans", "P1"],
        ["Events", "Ticketing take-rate + sponsorship", "After V45+ stable"],
        ["Mentorship", "FUTURE role — session commission later", "P3"],
        ["Partnerships", "CSR, NGO, city, campus, institutions", "Parallel P1"],
    ], rh=0.48, fs=12)

    s = keep(content(prs, "33  COMMERCIALIZATION  ·  2/2", "Long-term model", "RECOMMENDATION"))
    card(s, 0.4, 1.22, 6.2, 5.7, "North star economics",
         "Gross merchandise on Glow + doctors + jobs + products in one metro, at a transparent take-rate, plus provider boosters, plus 1–2 CSR/city contracts. Subscription is optional frosting. Panic is a public-good layer that earns trust, not INR. If SOS fails once in public, the take-rate is worthless.", ROSE)
    card(s, 6.8, 1.22, 6.1, 5.7, "What we will not do",
         "Paywall SOS. Sell SOS location to advertisers. Originate loans or run a fund. Pretend to be 112. Advertise ChatGPT features that are not in src/main/java. Expand to 10 cities before one city is dense and drilled.", NAVY)

    # close
    s = keep(content(prs, "CLOSE", "What a new stakeholder should remember", "FACT + RECOMMENDATION"))
    bullets(s, [
        "FightDFear is already a super-platform in code — JSP + Flutter + admin queues + Razorpay — not a concept note.",
        "It is not 112, not a bank, not an LLM product, not Thymeleaf, not a split RBAC system.",
        "Biggest 90-day work: schema, secrets, SOS actually notifies, audit log, dummy wallet, DPDP, legal lines.",
        "Market 2026: NCRB VAW still high; urban female CWS LFPR 24.8%; Life360 shows safety-graphs can be huge businesses; Naukri shows jobs cash; Sitha/SHE-Marts show women-economy is the political + commercial wave.",
        "Win by being the trusted local OS that 112 cannot be and Naukri will not be — in one city first.",
        "Companion narrative: docs/FightDFear_Product_Technical_Blueprint.md",
    ], size=15)

    total = len(built)
    for i, sl in enumerate(built, 1):
        # cover already branded
        if i == 1:
            continue
        # dividers have their own bar
        # still add page number on dividers via small box
        has_navy_full = False
        footer(sl, i, total)

    out = Path(r"C:\Users\Hp\Documents\KishorDfire\docs\FightDFear_33_Section_Blueprint.pptx")
    prs.save(str(out))
    print(f"Wrote {out} slides={total}")
    return out, total


if __name__ == "__main__":
    build()
